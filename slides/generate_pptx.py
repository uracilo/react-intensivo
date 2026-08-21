#!/usr/bin/env python3
"""Genera PPTX simples a partir de los Markdown en slides/."""

from __future__ import annotations

import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise SystemExit(
        "Instalá python-pptx: pip install python-pptx\n" + str(exc)
    ) from exc

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "pptx"
OUT.mkdir(exist_ok=True)

FILES = [
    "00-web-basica.md",
    "01-fundamentos.md",
    "02-estado.md",
    "03-datos.md",
    "04-crud-router.md",
    "05-taskflow-jwt.md",
]


def parse_slides(md: str) -> list[tuple[str, list[str]]]:
    chunks = re.split(r"\n---\n", md.strip())
    slides: list[tuple[str, list[str]]] = []
    for chunk in chunks:
        lines = [ln.rstrip() for ln in chunk.strip().splitlines() if ln.strip()]
        if not lines:
            continue
        title = lines[0].lstrip("# ").strip()
        body: list[str] = []
        for ln in lines[1:]:
            if ln.startswith("```"):
                continue
            body.append(ln.lstrip("- ").strip())
        slides.append((title, body))
    return slides


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x1F)

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    p.font.name = "Calibri"

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.8), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for line in bullets:
        if not line:
            continue
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = line
        para.level = 0
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(0xE7, 0xE5, 0xE4)
        para.font.name = "Calibri"
        para.space_after = Pt(10)


def build_deck(md_path: Path) -> Path:
    slides = parse_slides(md_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for title, body in slides:
        add_slide(prs, title, body)
    out = OUT / (md_path.stem + ".pptx")
    prs.save(out)
    return out


def main() -> None:
    for name in FILES:
        path = ROOT / name
        if not path.exists():
            print("skip missing", name)
            continue
        out = build_deck(path)
        print("wrote", out.name)


if __name__ == "__main__":
    main()
